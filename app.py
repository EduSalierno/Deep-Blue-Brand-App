import streamlit as st
import openai
import os
from PyPDF2 import PdfReader
from ebooklib import epub
from bs4 import BeautifulSoup
import docx
from pptx import Presentation
import pandas as pd
import re
import networkx as nx
import matplotlib.pyplot as plt
from io import BytesIO
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.cluster import KMeans
from msal_streamlit_authentication import msal_authentication
import requests

# MSAL Config
CLIENT_ID = st.secrets["ms_client_id"]
CLIENT_SECRET = st.secrets["ms_client_secret"]
TENANT_ID = st.secrets["ms_tenant_id"]
AUTHORITY = f"https://login.microsoftonline.com/{TENANT_ID}"
SCOPES = ["User.Read", "Files.ReadWrite.All", "offline_access"]

st.set_page_config(page_title="Deep Blue Brand System", layout="wide")
st.title("🔵 Deep Blue Brand System")
st.subheader("Analisi semiotica e strategia di brand potenziata da AI")

# Login OneDrive
result = msal_authentication(
    auth_url=AUTHORITY,
    client_id=CLIENT_ID,
    client_secret=CLIENT_SECRET,
    scopes=SCOPES,
)

if result:
    access_token = result["access_token"]
    user_info = requests.get(
        "https://graph.microsoft.com/v1.0/me",
        headers={"Authorization": f"Bearer {access_token}"}
    ).json()
    st.success(f"Autenticato come {user_info['displayName']}")

    # Caricamento documenti
    st.sidebar.header("📂 Caricamento documenti")
    uploaded_file = st.sidebar.file_uploader("Carica un file (PDF, EPUB, Word, PowerPoint, Excel)", type=["pdf", "epub", "docx", "pptx", "xlsx", "xls"])

    full_text = ""
    if uploaded_file:
        st.sidebar.success("File caricato correttamente!")
        file_extension = os.path.splitext(uploaded_file.name)[1].lower()

        if file_extension == ".pdf":
            reader = PdfReader(uploaded_file)
            full_text = "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])
        elif file_extension == ".epub":
            book = epub.read_epub(uploaded_file)
            chapters = []
            for item in book.get_items():
                if item.get_type() == epub.ITEM_DOCUMENT:
                    soup = BeautifulSoup(item.get_content(), 'html.parser')
                    chapters.append(soup.get_text())
            full_text = "\n".join(chapters)
        elif file_extension == ".docx":
            doc = docx.Document(uploaded_file)
            full_text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        elif file_extension == ".pptx":
            prs = Presentation(uploaded_file)
            slides_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slides_text.append(shape.text)
            full_text = "\n".join(slides_text)
        elif file_extension in [".xlsx", ".xls"]:
            df = pd.read_excel(uploaded_file)
            full_text = df.to_string(index=False)

        # Moduli di analisi
        st.sidebar.header("📊 Moduli di analisi")
        analysis_type = st.sidebar.selectbox("Scegli un'analisi", [
            "Analisi semiotica di base",
            "Atlante semiotico",
            "Insight strategici",
            "Canvas: Brand Core"
        ])

        if st.sidebar.button("Genera output"):
            with st.spinner("Analisi in corso..."):
                openai.api_key = st.secrets["openai_api_key"]
                prompt = f"Esegui una {analysis_type.lower()} sul seguente testo:\n\n{full_text}\n\nRispondi in modo strutturato, identificando codici, tropi, opposizioni binarie, valori culturali e insight strategici. Elenca i codici principali come lista."

                response = openai.ChatCompletion.create(
                    model="gpt-4",
                    messages=[
                        {"role": "system", "content": "Sei un esperto di semiotica e brand strategy."},
                        {"role": "user", "content": prompt}
                    ],
                    temperature=0.7
                )

                output_text = response.choices[0].message.content
                st.subheader("🧠 Risultato dell'analisi")
                st.markdown(output_text)
                st.download_button("📥 Scarica il risultato", output_text, file_name="analisi_deepblue.txt")

                requests.put(
                    url="https://graph.microsoft.com/v1.0/me/drive/root:/DeepBlue/analisi_deepblue.txt:/content",
                    headers={"Authorization": f"Bearer {access_token}", "Content-Type": "text/plain"},
                    data=output_text.encode("utf-8")
                )

                # Mappa visuale dei codici
                st.subheader("🕸️ Mappa visuale dei codici")
                code_matches = re.findall(r"- (.+?)\n", output_text)
                unique_codes = list(set(code_matches))

                if len(unique_codes) > 1:
                    G = nx.Graph()
                    for code in unique_codes:
                        G.add_node(code)
                    for i in range(len(unique_codes)):
                        for j in range(i + 1, len(unique_codes)):
                            G.add_edge(unique_codes[i], unique_codes[j])
                    fig, ax = plt.subplots(figsize=(10, 6))
                    pos = nx.spring_layout(G, seed=42)
                    nx.draw(G, pos, with_labels=True, node_color='lightblue', edge_color='gray', node_size=2000, font_size=10, ax=ax)
                    st.pyplot(fig)

                # Cluster semiotici
                st.subheader("🔍 Cluster semiotici tematici")
                try:
                    vectorizer = TfidfVectorizer(stop_words='english')
                    X = vectorizer.fit_transform(unique_codes)
                    num_clusters = min(5, len(unique_codes))
                    kmeans = KMeans(n_clusters=num_clusters, random_state=42).fit(X)
                    clusters = {i: [] for i in range(num_clusters)}
                    for i, label in enumerate(kmeans.labels_):
                        clusters[label].append(unique_codes[i])
                    for label, codes in clusters.items():
                        st.markdown(f"**Cluster {label+1}**: {', '.join(codes)}")
                except Exception as e:
                    st.warning(f"Impossibile generare cluster: {e}")
else:
    st.info("Effettua il login con Microsoft per iniziare.")
